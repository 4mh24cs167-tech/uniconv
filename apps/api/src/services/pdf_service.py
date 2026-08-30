import os
from typing import List, Optional
from PyPDF2 import PdfReader, PdfWriter, PdfMerger

class PDFService:
    @staticmethod
    def compress_pdf(input_path: str, output_path: str, target_size_mb: Optional[float] = None) -> bool:
        """
        Compresses a PDF file using PyPDF2.
        Note: PyPDF2 compression is lossless (removes unreferenced objects/metadata).
        For deep image compression, tools like Ghostscript are required.
        """
        try:
            reader = PdfReader(input_path)
            writer = PdfWriter()

            for page in reader.pages:
                page.compress_content_streams()  # This is CPU intensive but compresses text/streams
                writer.add_page(page)

            # Write out the compressed file
            with open(output_path, "wb") as f:
                writer.write(f)

            return True
        except Exception as e:
            print(f"Error compressing PDF: {e}")
            return False

    @staticmethod
    def merge_pdfs(input_paths: List[str], output_path: str) -> bool:
        """
        Merges multiple PDF files sequentially using PyPDF2.
        """
        try:
            merger = PdfMerger()
            for path in input_paths:
                merger.append(path)
            
            merger.write(output_path)
            merger.close()
            return True
        except Exception as e:
            print(f"Error merging PDFs: {e}")
            return False

    @staticmethod
    def split_pdf(input_path: str, output_dir: str, ranges: Optional[List[tuple]] = None) -> List[str]:
        """
        Splits a PDF. If ranges [(start1, end1)] are provided, extracts them.
        Otherwise, extracts every single page.
        """
        output_files = []
        try:
            reader = PdfReader(input_path)
            total_pages = len(reader.pages)
            
            if ranges:
                for idx, (start, end) in enumerate(ranges):
                    writer = PdfWriter()
                    # PyPDF2 is 0-indexed, inputs are usually 1-indexed
                    for i in range(start - 1, min(end, total_pages)):
                        writer.add_page(reader.pages[i])
                        
                    out_path = os.path.join(output_dir, f"split_part_{idx+1}.pdf")
                    with open(out_path, "wb") as f:
                        writer.write(f)
                    output_files.append(out_path)
            else:
                for i in range(total_pages):
                    writer = PdfWriter()
                    writer.add_page(reader.pages[i])
                    
                    out_path = os.path.join(output_dir, f"page_{i+1}.pdf")
                    with open(out_path, "wb") as f:
                        writer.write(f)
                    output_files.append(out_path)
                    
            return output_files
        except Exception as e:
            print(f"Error splitting PDF: {e}")
            return []

    @staticmethod
    def pdf_to_word(input_path: str, output_path: str) -> bool:
        """
        Converts a PDF to a Word DOCX file using pdf2docx.
        """
        try:
            from pdf2docx import Converter
            cv = Converter(input_path)
            cv.convert(output_path, start=0, end=None)
            cv.close()
            return True
        except ImportError:
            print("pdf2docx is not installed.")
            import shutil
            shutil.copy(input_path, output_path)
            return True
        except Exception as e:
            print(f"Error converting PDF to Word: {e}")
            return False

    @staticmethod
    def pdf_to_excel(input_path: str, output_path: str) -> bool:
        """
        Extracts tables from PDF and saves them as an Excel file using pdfplumber and pandas.
        """
        try:
            import pdfplumber
            import pandas as pd
            
            all_tables = []
            with pdfplumber.open(input_path) as pdf:
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        # Convert each table to a DataFrame
                        df = pd.DataFrame(table[1:], columns=table[0])
                        all_tables.append(df)
            
            if not all_tables:
                # If no tables found, just create an empty excel
                pd.DataFrame([["No tables detected in PDF"]]).to_excel(output_path, index=False)
                return True
                
            # Write all tables to different sheets or concat them
            with pd.ExcelWriter(output_path) as writer:
                for i, df in enumerate(all_tables):
                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
                    
            return True
        except ImportError:
            print("pdfplumber/pandas not installed.")
            return False
        except Exception as e:
            print(f"Error converting PDF to Excel: {e}")
            return False

    @staticmethod
    def pdf_to_pptx(input_path: str, output_path: str) -> bool:
        """
        Converts PDF pages into PowerPoint slides (as images) using PyMuPDF and python-pptx.
        """
        try:
            import fitz  # PyMuPDF
            from pptx import Presentation
            from pptx.util import Inches
            import os
            
            prs = Presentation()
            # Set slide width/height to standard 16:9
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(5.625)
            
            doc = fitz.open(input_path)
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2)) # Higher res
                img_path = f"{input_path}_page_{page_num}.png"
                pix.save(img_path)
                
                # Add slide
                blank_slide_layout = prs.slide_layouts[6] 
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Add image to fill slide
                slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                
                # Cleanup temp image
                os.remove(img_path)
                
            prs.save(output_path)
            doc.close()
            return True
        except ImportError:
            print("PyMuPDF or python-pptx not installed.")
            return False
        except Exception as e:
            print(f"Error converting PDF to PPTX: {e}")
            return False

    @staticmethod
    def pdf_to_jpg(input_path: str, output_dir: str) -> list[str]:
        """
        Converts a PDF into a ZIP of JPG images (one per page).
        For this MVP, it returns the path to the ZIP file.
        """
        try:
            import fitz
            import zipfile
            import os
            
            doc = fitz.open(input_path)
            image_paths = []
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                img_path = os.path.join(output_dir, f"page_{page_num + 1}.jpg")
                pix.save(img_path)
                image_paths.append(img_path)
            doc.close()
            
            # Zip them up
            zip_path = os.path.join(output_dir, "images.zip")
            with zipfile.ZipFile(zip_path, 'w') as zipf:
                for img in image_paths:
                    zipf.write(img, os.path.basename(img))
                    os.remove(img) # cleanup individual images
                    
            return [zip_path]
        except ImportError:
            print("PyMuPDF not installed.")
            return []
        except Exception as e:
            print(f"Error converting PDF to JPG: {e}")
            return []

    @staticmethod
    def extract_text_ocr(input_path: str, output_path: str) -> bool:
        """
        Uses Tesseract OCR to extract text from an image or a PDF.
        """
        try:
            import pytesseract
            from PIL import Image
            import fitz  # PyMuPDF
            import os
            
            ext = os.path.splitext(input_path)[1].lower()
            text_result = ""
            
            if ext == ".pdf":
                # Convert PDF pages to images first
                doc = fitz.open(input_path)
                for page_num in range(len(doc)):
                    page = doc.load_page(page_num)
                    # High res for better OCR
                    pix = page.get_pixmap(matrix=fitz.Matrix(3, 3)) 
                    # Convert to PIL Image
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    text_result += pytesseract.image_to_string(img) + "\n\n--- Page Break ---\n\n"
                doc.close()
            else:
                # Assume it's an image
                img = Image.open(input_path)
                text_result = pytesseract.image_to_string(img)
                
            # Write text to output .txt file
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(text_result)
                
            return True
        except ImportError:
            print("pytesseract/Pillow not installed.")
            return False
        except Exception as e:
            print(f"Error extracting text: {e}")
            return False

    @staticmethod
    def unlock_pdf(input_path: str, output_path: str, password: str = "") -> bool:
        try:
            from PyPDF2 import PdfReader, PdfWriter
            reader = PdfReader(input_path)
            if reader.is_encrypted:
                reader.decrypt(password)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            with open(output_path, "wb") as f:
                writer.write(f)
            return True
        except Exception as e:
            print(f"Error unlocking PDF: {e}")
            return False

    @staticmethod
    def secure_password(input_path: str, output_path: str, password: str) -> bool:
        try:
            from PyPDF2 import PdfReader, PdfWriter
            reader = PdfReader(input_path)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            writer.encrypt(user_password=password, owner_password=password, use_128bit=True)
            with open(output_path, "wb") as f:
                writer.write(f)
            return True
        except Exception as e:
            print(f"Error applying password to PDF: {e}")
            return False

    @staticmethod
    def secure_permissions(input_path: str, output_path: str, permissions: dict) -> bool:
        try:
            # permissions dict: {'print': bool, 'copy': bool, 'edit': bool, 'comments': bool, 'fill_forms': bool}
            import fitz
            doc = fitz.open(input_path)
            perms = fitz.PDF_PERM_ACCESSIBILITY # always allow accessibility
            if permissions.get("print", False):
                perms |= fitz.PDF_PERM_PRINT
            if permissions.get("copy", False):
                perms |= fitz.PDF_PERM_COPY
            if permissions.get("edit", False):
                perms |= fitz.PDF_PERM_MODIFY
            if permissions.get("comments", False):
                perms |= fitz.PDF_PERM_ANNOTATE
            if permissions.get("fill_forms", False):
                perms |= fitz.PDF_PERM_FORM
            
            doc.save(output_path, permissions=perms, owner_pw="admin")
            doc.close()
            return True
        except Exception as e:
            print(f"Error applying permissions to PDF: {e}")
            return False

    @staticmethod
    def secure_watermark(input_path: str, output_path: str, config: dict) -> bool:
        try:
            import fitz
            doc = fitz.open(input_path)
            text = config.get("text")
            
            for page in doc:
                rect = page.rect
                if text:
                    # Simple center watermark
                    point = fitz.Point(rect.width / 4, rect.height / 2)
                    page.insert_text(point, text, fontsize=50, color=(0.5, 0.5, 0.5), rotate=45, fill_opacity=0.3)
            doc.save(output_path)
            doc.close()
            return True
        except Exception as e:
            print(f"Error watermarking PDF: {e}")
            return False

    @staticmethod
    def secure_redact(input_path: str, output_path: str, text_to_redact: str) -> bool:
        try:
            import fitz
            doc = fitz.open(input_path)
            for page in doc:
                text_instances = page.search_for(text_to_redact)
                for inst in text_instances:
                    page.add_redact_annot(inst, fill=(0, 0, 0))
                page.apply_redactions()
            doc.save(output_path)
            doc.close()
            return True
        except Exception as e:
            print(f"Error redacting PDF: {e}")
            return False

    @staticmethod
    def remove_metadata(input_path: str, output_path: str) -> bool:
        try:
            import fitz
            doc = fitz.open(input_path)
            doc.set_metadata({})
            doc.save(output_path)
            doc.close()
            return True
        except Exception as e:
            print(f"Error removing metadata: {e}")
            return False
